# This module handles the ingestion of various document types, extracting text content, splitting it into chunks, and storing it in a vector database for later retrieval by the agent.

import logging
import uuid
import io
import os
import csv
import json
import asyncio
import tempfile

import pypdf
import docx
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load as odf_load
from odf.text import P
import ebooklib
from ebooklib import epub

from langchain_text_splitters import RecursiveCharacterTextSplitter
from app.services.vector_store import get_vector_store
from app.database import get_db_connection
from app.services.ocr import ocr_pdf

logger = logging.getLogger(__name__)
USE_POSTGRES = os.getenv("USE_POSTGRES", "false").lower() == "true"

class DocumentProcessor:
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", " ", ""]
        )
        self.vector_store = get_vector_store()

    def _extract_text(self, content: bytes, filename: str) -> str:
        """Routes file content bytes to the appropriate parsing library based on extension."""
        ext = filename.lower().split('.')[-1]
        
        if ext in ['txt', 'md']:
            return content.decode("utf-8", errors="ignore")
            
        
        elif ext == 'pdf':
            # 1. First, try standard PyPDF extraction (fast, works for text-based PDFs)
            pdf_reader = pypdf.PdfReader(io.BytesIO(content))
            extracted = "".join((page.extract_text() or "") + "\n" for page in pdf_reader.pages)
            
            # 2. If extracted text is very short (likely scanned) and OCR is enabled, fall back to OCR
            if len(extracted.strip()) < 100 and os.getenv("ENABLE_OCR", "false").lower() == "true":
                try:
                    ocr_text = ocr_pdf(content)
                    if ocr_text and len(ocr_text.strip()) > 0:
                        logger.info(f"OCR succeeded for {filename} (chars: {len(ocr_text)})")
                        return ocr_text
                    else:
                        logger.warning(f"OCR returned empty for {filename}, keeping PyPDF extraction.")
                except Exception as e:
                    logger.error(f"OCR exception for {filename}: {e}, keeping PyPDF extraction.")
            
            # 3. Return the original PyPDF extraction (may be empty for scanned PDFs)
            return extracted
            
        elif ext == 'docx':
            doc = docx.Document(io.BytesIO(content))
            return "\n".join([p.text for p in doc.paragraphs])
            
        elif ext == 'pptx':
            prs = Presentation(io.BytesIO(content))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
            
        elif ext == 'xlsx':
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text.append(row_text)
            return "\n".join(text)
            
        elif ext == 'csv':
            decoded = content.decode("utf-8", errors="ignore")
            reader = csv.reader(io.StringIO(decoded))
            return "\n".join([" ".join(row) for row in reader])
            
        elif ext == 'json':
            data = json.loads(content.decode("utf-8", errors="ignore"))
            return json.dumps(data, indent=2)
            
        elif ext in ['html', 'xml']:
            soup = BeautifulSoup(content.decode("utf-8", errors="ignore"), "html.parser")
            return soup.get_text(separator="\n", strip=True)
            
        elif ext == 'rtf':
            return rtf_to_text(content.decode("utf-8", errors="ignore"))
            
        elif ext == 'odt':
            doc = odf_load(io.BytesIO(content))
            return "\n".join([str(p) for p in doc.getElementsByType(P)])
            
        elif ext == 'epub':
            with tempfile.NamedTemporaryFile(delete=False, suffix=".epub") as temp:
                temp.write(content)
                temp_path = temp.name
            
            try:
                book = epub.read_epub(temp_path)
                text = []
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_body_content(), 'html.parser')
                        text.append(soup.get_text(strip=True))
                return "\n".join(text)
            finally:
                os.remove(temp_path)
                
        else:
            raise ValueError(f"No text extraction logic defined for extension: {ext}")

    def _split_and_embed(self, text_content: str, filename: str, owner_id: str) -> int:
        chunks = self.splitter.split_text(text_content)
        if not chunks:
            return 0
            
        metadatas = [{"source": filename, "chunk_index": i, "owner_id": owner_id} for i in range(len(chunks))]
        ids = [str(uuid.uuid4()) for _ in range(len(chunks))]
        self.vector_store.add_texts(texts=chunks, metadatas=metadatas, ids=ids)
        return len(chunks)

    def _insert_metadata(self, filename: str, file_type: str, chunk_count: int, owner_id: str) -> int:
        with get_db_connection() as conn:
            if USE_POSTGRES:
                with conn.cursor() as cur:
                    cur.execute(
                        "INSERT INTO documents (filename, file_type, chunk_count, owner_id) VALUES (%s, %s, %s, %s) RETURNING id",
                        (filename, file_type, chunk_count, owner_id)
                    )
                    doc_id = cur.fetchone()[0]
            else:
                cur = conn.execute(
                    "INSERT INTO documents (filename, file_type, chunk_count, owner_id) VALUES (?, ?, ?, ?)",
                    (filename, file_type, chunk_count, owner_id)
                )
                doc_id = cur.lastrowid
            conn.commit()
            return doc_id

    def _update_chunk_count(self, doc_id: int, chunk_count: int):
        with get_db_connection() as conn:
            if USE_POSTGRES:
                with conn.cursor() as cur:
                    cur.execute("UPDATE documents SET chunk_count = %s WHERE id = %s", (chunk_count, doc_id))
            else:
                conn.execute("UPDATE documents SET chunk_count = ? WHERE id = ?", (chunk_count, doc_id))
            conn.commit()

    async def process_and_store(self, filename: str, content: bytes, owner_id: str) -> None:
        logger.info(f"Starting async ingestion for {filename} for owner {owner_id}")
        try:
            file_type = filename.split('.')[-1].lower()
            # 1. Insert metadata immediately with chunk_count=0
            doc_id = await asyncio.to_thread(self._insert_metadata, filename, file_type, 0, owner_id)
            logger.info(f"Metadata inserted for {filename} (id: {doc_id})")

            # 2. Extract text
            # Extract text, with a timeout for OCR (if enabled)
            try:
                text_content = await asyncio.wait_for(
                    asyncio.to_thread(self._extract_text, content, filename),
                    timeout=60.0  # 60 seconds max for OCR; adjust as needed
                )
            except asyncio.TimeoutError:
                logger.warning(f"Text extraction timed out for {filename}, falling back to empty text.")
                text_content = ""  # or you could try PyPDF directly, but we'll just treat as empty

            if not text_content.strip():
                logger.warning(f"No extractable text found in {filename}.")
                await asyncio.to_thread(self._update_chunk_count, doc_id, -1)
                return

            # 3. Split and embed
            chunk_count = await asyncio.to_thread(self._split_and_embed, text_content, filename, owner_id)

            # 4. Update chunk count
            if chunk_count > 0:
                await asyncio.to_thread(self._update_chunk_count, doc_id, chunk_count)

            logger.info(f"Ingestion complete for {filename}. Chunks: {chunk_count}")

        except Exception as e:
            logger.error(f"Ingestion failed for {filename}: {str(e)}", exc_info=True)
            raise